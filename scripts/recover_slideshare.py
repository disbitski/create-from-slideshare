#!/usr/bin/env python3
"""Recover a visible Slideshare deck into images, PDF, and image-based PPTX."""

from __future__ import annotations

import argparse
import html
import json
import re
import sys
import time
import urllib.error
import urllib.parse
import urllib.request
import zipfile
from dataclasses import dataclass
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches


USER_AGENT = (
    "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
    "AppleWebKit/537.36 (KHTML, like Gecko) Chrome/126 Safari/537.36"
)

SLIDE_URL_RE = re.compile(
    r"https://image\.slidesharecdn\.com/"
    r"(?P<deck>[^\"'<>\\\s]+?)/(?P<tier>\d+)/"
    r"(?P<stem>[^\"'<>\\\s]+?)-(?P<num>\d+)-(?P<size>\d+)"
    r"\.(?P<ext>jpg|jpeg|png|webp)",
    re.IGNORECASE,
)


@dataclass(frozen=True)
class SlideCandidate:
    slide: int
    url: str
    observed_size: int


def request_bytes(url: str, method: str = "GET") -> tuple[bytes, str]:
    request = urllib.request.Request(url, method=method, headers={"User-Agent": USER_AGENT})
    with urllib.request.urlopen(request, timeout=30) as response:
        return response.read(), response.headers.get("content-type", "")


def normalize_html(raw: str) -> str:
    normalized = html.unescape(raw)
    normalized = normalized.replace("\\u0026", "&")
    normalized = normalized.replace("\\/", "/")
    return normalized


def fetch_page(url: str) -> str:
    data, _ = request_bytes(url)
    return normalize_html(data.decode("utf-8", errors="replace"))


def extract_candidates(page_html: str) -> dict[int, list[SlideCandidate]]:
    candidates: dict[int, list[SlideCandidate]] = {}
    seen: set[tuple[int, str]] = set()

    for match in SLIDE_URL_RE.finditer(page_html):
        slide = int(match.group("num"))
        observed_size = int(match.group("size"))
        base = (
            f"https://image.slidesharecdn.com/{match.group('deck')}/"
            f"{match.group('tier')}/{match.group('stem')}-{slide}"
        )
        url = match.group(0)
        variants = [
            f"https://image.slidesharecdn.com/{match.group('deck')}/75/{match.group('stem')}-{slide}-2048.jpg",
            f"{base}-2048.jpg",
            f"https://image.slidesharecdn.com/{match.group('deck')}/95/{match.group('stem')}-{slide}-1024.jpg",
            f"{base}-1024.jpg",
            url,
        ]

        for variant in variants:
            key = (slide, variant)
            if key in seen:
                continue
            seen.add(key)
            candidates.setdefault(slide, []).append(
                SlideCandidate(slide=slide, url=variant, observed_size=observed_size)
            )

    return dict(sorted(candidates.items()))


def extension_for_content_type(content_type: str, fallback_url: str) -> str:
    content_type = content_type.lower().split(";")[0].strip()
    if content_type == "image/webp":
        return ".webp"
    if content_type in {"image/jpeg", "image/jpg"}:
        return ".jpg"
    if content_type == "image/png":
        return ".png"
    suffix = Path(urllib.parse.urlparse(fallback_url).path).suffix.lower()
    return suffix if suffix in {".webp", ".jpg", ".jpeg", ".png"} else ".img"


def download_slides(
    candidates_by_slide: dict[int, list[SlideCandidate]],
    raw_dir: Path,
    expected_slides: int | None,
) -> list[dict]:
    if not candidates_by_slide:
        raise RuntimeError("No Slideshare slide image URLs were found in the page HTML.")

    slide_numbers = sorted(candidates_by_slide)
    if expected_slides is not None and len(slide_numbers) != expected_slides:
        raise RuntimeError(
            f"Expected {expected_slides} slides but found {len(slide_numbers)} slide numbers."
        )

    raw_dir.mkdir(parents=True, exist_ok=True)
    manifest = []

    for slide in slide_numbers:
        last_error = None
        for candidate in candidates_by_slide[slide]:
            try:
                data, content_type = request_bytes(candidate.url)
                if not content_type.lower().startswith("image/") or len(data) < 1000:
                    raise RuntimeError(f"unexpected response {content_type!r}, {len(data)} bytes")
                extension = extension_for_content_type(content_type, candidate.url)
                filename = f"slide-{slide:03d}{extension}"
                destination = raw_dir / filename
                destination.write_bytes(data)
                manifest.append(
                    {
                        "slide": slide,
                        "url": candidate.url,
                        "raw_file": str(destination),
                        "content_type": content_type,
                        "bytes": len(data),
                    }
                )
                print(f"downloaded slide {slide:03d} from {candidate.url}")
                break
            except (urllib.error.URLError, TimeoutError, RuntimeError) as error:
                last_error = error
                time.sleep(0.15)
        else:
            raise RuntimeError(f"Could not download slide {slide}: {last_error}")

    return manifest


def convert_to_jpegs(manifest: list[dict], image_dir: Path) -> list[Path]:
    image_dir.mkdir(parents=True, exist_ok=True)
    output_paths = []

    for item in manifest:
        slide = item["slide"]
        source = Path(item["raw_file"])
        destination = image_dir / f"slide-{slide:03d}.jpg"
        with Image.open(source) as image:
            rgb = image.convert("RGB")
            rgb.save(destination, "JPEG", quality=94, optimize=True)
            item["width"] = rgb.width
            item["height"] = rgb.height
            item["image_file"] = str(destination)
        output_paths.append(destination)
        print(f"converted slide {slide:03d} to JPEG")

    return output_paths


def build_pdf(image_paths: list[Path], pdf_path: Path) -> None:
    pages = [Image.open(path).convert("RGB") for path in image_paths]
    first, rest = pages[0], pages[1:]
    first.save(pdf_path, "PDF", save_all=True, append_images=rest, resolution=144.0)
    for page in pages:
        page.close()


def build_pptx(image_paths: list[Path], pptx_path: Path) -> None:
    with Image.open(image_paths[0]) as first:
        width, height = first.size

    deck = Presentation()
    deck.slide_width = Inches(10)
    deck.slide_height = Inches(10 * height / width)
    blank_layout = deck.slide_layouts[6]

    for image_path in image_paths:
        slide = deck.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(image_path),
            0,
            0,
            width=deck.slide_width,
            height=deck.slide_height,
        )

    deck.save(pptx_path)


def build_contact_sheet(image_paths: list[Path], contact_sheet_path: Path) -> None:
    thumb_w, thumb_h, label_h = 320, 240, 28
    cols = 6
    rows = (len(image_paths) + cols - 1) // cols
    sheet = Image.new("RGB", (cols * thumb_w, rows * (thumb_h + label_h)), "white")
    draw = ImageDraw.Draw(sheet)

    for index, image_path in enumerate(image_paths, start=1):
        row, col = divmod(index - 1, cols)
        x = col * thumb_w
        y = row * (thumb_h + label_h)
        with Image.open(image_path) as image:
            thumbnail = image.convert("RGB")
            thumbnail.thumbnail((thumb_w, thumb_h), Image.LANCZOS)
            sheet.paste(thumbnail, (x + (thumb_w - thumbnail.width) // 2, y))
        draw.text((x + 8, y + thumb_h + 6), f"Slide {index}", fill=(0, 0, 0))

    sheet.save(contact_sheet_path, "JPEG", quality=90)


def pptx_slide_count(pptx_path: Path) -> int:
    with zipfile.ZipFile(pptx_path) as archive:
        return sum(
            1
            for name in archive.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        )


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("url", help="Slideshare slideshow URL")
    parser.add_argument("--output-dir", required=True, help="Directory for recovered files")
    parser.add_argument("--expected-slides", type=int, help="Fail if the page exposes a different count")
    parser.add_argument(
        "--basename",
        default="deck-recovered-from-slideshare",
        help="Base filename for PPTX and PDF outputs",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    output_dir = Path(args.output_dir).expanduser().resolve()
    raw_dir = output_dir / "raw"
    image_dir = output_dir / "images"
    output_dir.mkdir(parents=True, exist_ok=True)

    page = fetch_page(args.url)
    candidates = extract_candidates(page)
    manifest = download_slides(candidates, raw_dir, args.expected_slides)
    image_paths = convert_to_jpegs(manifest, image_dir)

    pdf_path = output_dir / f"{args.basename}.pdf"
    pptx_path = output_dir / f"{args.basename}.pptx"
    contact_sheet_path = output_dir / "contact-sheet.jpg"
    manifest_path = output_dir / "manifest.json"

    build_pdf(image_paths, pdf_path)
    build_pptx(image_paths, pptx_path)
    build_contact_sheet(image_paths, contact_sheet_path)

    result = {
        "source_url": args.url,
        "slide_count": len(image_paths),
        "pptx_slide_count": pptx_slide_count(pptx_path),
        "pptx": str(pptx_path),
        "pdf": str(pdf_path),
        "contact_sheet": str(contact_sheet_path),
        "images_dir": str(image_dir),
        "raw_dir": str(raw_dir),
        "slides": manifest,
    }
    manifest_path.write_text(json.dumps(result, indent=2), encoding="utf-8")
    print(json.dumps(result, indent=2))
    return 0


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except Exception as exc:  # noqa: BLE001
        print(f"ERROR: {exc}", file=sys.stderr)
        raise SystemExit(1)
